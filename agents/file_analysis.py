# analyst_ingest_and_summarize.py
# Handles: PDF, DOCX, PPTX -> Structured Analysis -> Firestore or Local

import os, sys, json, uuid, datetime, tempfile, requests
import pdfplumber
from docx import Document
from pptx import Presentation
from dotenv import load_dotenv
from pathlib import Path
from typing import Dict, Any, List

import firebase_admin
from firebase_admin import credentials, firestore
from google import genai

# ---------------------- Config ----------------------

DEFAULT_MODEL = os.getenv("GENAI_MODEL", "gemini-2.0-flash-lite-001")
DEFAULT_PROJECT = os.getenv("GENAI_PROJECT", "test-project-471516")
DEFAULT_LOCATION = os.getenv("GENAI_LOCATION", "us-central1")

GENERATION = {
    "temperature": 0.2,
    "top_p": 0.9,
    "max_output_tokens": 800,
    "candidate_count": 1
}

# ---------------------- Firebase Init ----------------------

if not firebase_admin._apps:
    cred_path = "startupiq-c1fe1-firebase-adminsdk-fbsvc-8b23336991.json"
    if os.path.exists(cred_path):
        cred = credentials.Certificate(cred_path)
        firebase_admin.initialize_app(cred)

db = firestore.client() if firebase_admin._apps else None

# ---------------------- Utilities ----------------------

def download_files(file_links: List[str]) -> List[str]:
    """Download files from URLs into local temp paths."""
    local_paths = []
    for url in file_links:
        try:
            r = requests.get(url, timeout=30)
            if r.status_code == 200:
                suffix = url.split("?")[0].split(".")[-1].lower()
                tmp = tempfile.NamedTemporaryFile(delete=False, suffix=f".{suffix}")
                tmp.write(r.content)
                tmp.close()
                local_paths.append(tmp.name)
            else:
                print(f"⚠️ Failed to download {url} (status {r.status_code})")
        except Exception as e:
            print(f"⚠️ Error downloading {url}: {e}")
    return local_paths

def extract_text_from_file(path: str) -> str:
    """Extract text depending on file type."""
    try:
        if path.endswith(".pdf"):
            text = ""
            with pdfplumber.open(path) as pdf:
                for page in pdf.pages:
                    text += (page.extract_text() or "") + "\n"
            return text.strip()

        elif path.endswith(".docx"):
            doc = Document(path)
            return "\n".join([p.text for p in doc.paragraphs if p.text.strip()])

        elif path.endswith(".pptx"):
            prs = Presentation(path)
            slides_text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slides_text.append(shape.text)
            return "\n".join(slides_text)

        else:
            return f"[Unsupported file format: {path}]"
    except Exception as e:
        return f"[Error extracting {path}: {e}]"

def extract_texts(file_paths: List[str]) -> Dict[str, str]:
    """Extract text from multiple files into dict {path: text}."""
    return {path: extract_text_from_file(path) for path in file_paths}

def extract_resp_text(resp) -> str:
    """Extract text from GenAI response object."""
    if hasattr(resp, "text") and isinstance(resp.text, str):
        return resp.text
    try:
        return resp.candidates[0].content.parts[0].text
    except Exception:
        return str(resp)

# ---------------------- Structured Analysis ----------------------

SCHEMA = {
    "idea_summary": {
        "problem": "",
        "solution": "",
        "alternatives": [],
        "target_market": "",
        "business_model": "",
        "go_to_market_strategy": ""
    },
    "company_details": {
        "name": "",
        "legal_entity": "",
        "sector": "",
        "sub_sector": "",
        "description": "",
        "headquarters": "",
        "employee_count": ""
    },
    "founders": [],
    "advisory_board": [],
    "partnerships": [],
    "financials_market": {
        "revenue_streams": [],
        "pricing_strategy": "",
        "unit_economics": {
            "cac": "",
            "ltv": "",
            "payback_period_months": ""
        },
        "traction": {
            "customers": "",
            "growth_rate": ""
        },
        "fundraising": {
            "stage": "",
            "amount_raised": "",
            "investors": []
        },
        "risks": []
    }
}

def merge_with_schema(base, parsed):
    """Ensure parsed JSON has all schema keys, fill empty if missing."""
    if isinstance(base, dict):
        out = {}
        for k, v in base.items():
            if k in parsed:
                out[k] = merge_with_schema(v, parsed[k])
            else:
                out[k] = v
        return out
    else:
        return parsed if parsed not in [None, ""] else base

def analyze_texts(texts: Dict[str, str], client, model=DEFAULT_MODEL) -> Dict[str, Any]:
    """Analyze document texts into structured JSON."""
    analysis = {}
    for path, content in texts.items():
        if not content.strip():
            analysis[path] = SCHEMA
            continue

        prompt = f"""
        You are an analyst. Extract the following structured JSON from the document. 
        Fill in values where possible. If something is missing, keep the key but use empty string, 0, or [].

        Required JSON schema:
        {json.dumps(SCHEMA, indent=2)}

        Document text:
        {content[:12000]}
        """

        try:
            resp = client.models.generate_content(
                model=model,
                contents=prompt,
                config={
                    "temperature": GENERATION["temperature"],
                    "max_output_tokens": GENERATION["max_output_tokens"],
                    "top_p": GENERATION["top_p"],
                }
            )
            text_resp = extract_resp_text(resp)

            parsed = None
            try:
                parsed = json.loads(text_resp)
            except Exception:
                print(f"⚠️ Could not parse JSON for {path}, storing empty schema")
                parsed = SCHEMA

            analysis[path] = merge_with_schema(SCHEMA, parsed)

        except Exception as e:
            print(f"⚠️ Error analyzing {path}: {e}")
            analysis[path] = SCHEMA

    return analysis

# ---------------------- Main Flow ----------------------

def process_firestore_doc(doc_id: str):
    print(f"Processing Firestore doc: {doc_id}")
    if not db:
        print("⚠️ Firebase not initialized")
        return
    doc = db.collection("files-map").document(doc_id).get()
    if not doc.exists:
        print("⚠️ Document not found in Firestore")
        return
    data = doc.to_dict()
    file_links = data.get("fileLinks", [])
    if not file_links:
        print("⚠️ No fileLinks found")
        return
    file_paths = download_files(file_links)
    texts = extract_texts(file_paths)
    client = genai.Client(vertexai=True, project=DEFAULT_PROJECT, location=DEFAULT_LOCATION,  credentials="test-project-471516-d93230f60a85.json")
    analysis = analyze_texts(texts, client)
    run_id = uuid.uuid4().hex
    ts = datetime.datetime.utcnow().isoformat() + "Z"
    db.collection("files-map").document(doc_id).update({
        "analysis": analysis,
        "analyzedAt": firestore.SERVER_TIMESTAMP,
        "analysis_run_id": run_id,
        "analysis_ts": ts
    })
    print(f"✅ Analysis stored for doc {doc_id}")

def process_local_files(file_paths: List[str]):
    texts = extract_texts(file_paths)
    client = genai.Client(vertexai=True, project=DEFAULT_PROJECT, location=DEFAULT_LOCATION, credentials="test-project-471516-d93230f60a85.json")
    analysis = analyze_texts(texts, client)
    print(json.dumps(analysis, indent=2))

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python analyst_ingest_and_summarize.py <doc_id|file1 file2 ...>")
        sys.exit(1)

    if os.path.exists(sys.argv[1]):
        # Treat as local file(s)
        files = sys.argv[1:]
        process_local_files(files)
    else:
        # Treat as Firestore doc_id
        doc_id = sys.argv[1]
        process_firestore_doc(doc_id)
